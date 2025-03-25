from pptx import Presentation
from pptx.util import Inches, Pt, RGBColor
from pptx.enum.text import PP_ALIGN

def add_textbox_with_dynamic_sizing(slide, left, top, width, height, text, set_border=True, min_font_size=6, max_font_size=24):
    """
    Add a text box with dynamically adjusted font size to fit the given dimensions.
    
    :param slide: PowerPoint slide to add the text box
    :param left: Left position of the text box
    :param top: Top position of the text box
    :param width: Width of the text box
    :param height: Height of the text box
    :param text: Text to be added
    :param set_border: Whether to add a border to the text box
    :param min_font_size: Minimum font size to prevent text from becoming unreadable
    :param max_font_size: Maximum font size to prevent overflow
    :return: Created text box
    """
    # Create the text box
    textbox = slide.shapes.add_textbox(left, top, width, height)
    
    # Add border if requested
    if set_border:
        textbox.line.color.rgb = RGBColor(0, 0, 0)  # black
        textbox.line.width = Pt(1)
    
    # Get the text frame
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # Function to check if text fits the text box
    def does_text_fit(font_size):
        # Reset the text frame
        tf.clear()
        
        # Process input text line by line
        for line in text.split('\n'):
            p = tf.add_paragraph()
            
            # Handle lines with colon (for bold/underlined first part)
            if ':' in line:
                parts = line.split(':', 1)
                run_bold = p.add_run()
                run_bold.text = parts[0] + ": "
                run_bold.font.bold = True
                run_bold.font.underline = True
                run_bold.font.name = 'Calibri'
                run_bold.font.size = Pt(font_size)
                
                run_normal = p.add_run()
                run_normal.text = parts[1]
                run_normal.font.name = 'Calibri'
                run_normal.font.size = Pt(font_size)
            else:
                p.text = line
                p.font.size = Pt(font_size)
        
        # Check if text overflows
        return tf.text_height <= height and tf.text_width <= width
    
    # Binary search to find optimal font size
    left_size, right_size = min_font_size, max_font_size
    best_font_size = max_font_size
    
    while left_size <= right_size:
        mid_size = (left_size + right_size) // 2
        
        if does_text_fit(mid_size):
            # If text fits, try to increase font size
            best_font_size = mid_size
            left_size = mid_size + 1
        else:
            # If text doesn't fit, reduce font size
            right_size = mid_size - 1
    
    # Reset text frame and add text with optimal font size
    tf.clear()
    for line in text.split('\n'):
        p = tf.add_paragraph()
        
        if ':' in line:
            parts = line.split(':', 1)
            run_bold = p.add_run()
            run_bold.text = parts[0] + ": "
            run_bold.font.bold = True
            run_bold.font.underline = True
            run_bold.font.name = 'Calibri'
            run_bold.font.size = Pt(best_font_size)
            
            run_normal = p.add_run()
            run_normal.text = parts[1]
            run_normal.font.name = 'Calibri'
            run_normal.font.size = Pt(best_font_size)
        else:
            p.text = line
            p.font.size = Pt(best_font_size)
    
    return textbox

# Example usage
def main():
    # Create a new presentation
    prs = Presentation()
    
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout
    
    # Add a text box with dynamic sizing
    text = "This is a long text that needs to be dynamically sized to fit the text box.\nIt can handle multiple lines and special formatting like: Bold Underlined Text"
    add_textbox_with_dynamic_sizing(
        slide, 
        left=Inches(1), 
        top=Inches(2), 
        width=Inches(6), 
        height=Inches(2), 
        text=text
    )
    
    # Save the presentation
    prs.save('dynamic_textbox_example.pptx')

if __name__ == '__main__':
    main()