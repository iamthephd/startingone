from io import BytesIO
from pptx import Presentation

def generate_ppt(commentary, selected_cells, file_name):
    """
    Generate PowerPoint presentation
    """
    try:
        prs = Presentation()
        
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = f"Analysis Report: {file_name}"
        
        # Add content slide
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        title.text = "Analysis Details"
        
        # Add commentary
        content = content_slide.placeholders[1]
        content.text = commentary
        
        # Save to BytesIO
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
    except Exception as e:
        print(f"Error generating PPT: {str(e)}")
        return BytesIO()  # Return empty buffer on error
